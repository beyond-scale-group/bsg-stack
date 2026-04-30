#!/usr/bin/env python3
"""
Render markdown to PPTX via python-pptx.

Slide structure conventions:
  # Title        -> title slide (layout 0: Title Slide)
  ## Section     -> content slide (layout 1: Title and Content)
  ### Sub        -> subheading within current slide
  ---            -> explicit slide break

Body content:
  - Bullets      -> text frame with bullet paragraphs
  | Tables |     -> PPTX table shape
  ![alt](path)  -> embedded picture
  **bold**       -> bold runs
  *italic*       -> italic runs
"""

from __future__ import annotations

import argparse
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path


# ---------------------------------------------------------------------------
# Markdown parser
# ---------------------------------------------------------------------------

@dataclass
class Slide:
    layout: str  # "title" | "content"
    title: str = ""
    blocks: list = field(default_factory=list)


def strip_frontmatter(text: str) -> str:
    if not text.startswith("---\n"):
        return text
    end = text.find("\n---\n", 4)
    if end != -1:
        return text[end + 5:]
    end = text.find("\n---", 4)
    if end != -1 and end + 4 >= len(text.rstrip()):
        return ""
    return text


def _parse_table(lines: list[str]) -> dict:
    rows = []
    for line in lines:
        cells = [c.strip() for c in line.strip().strip("|").split("|")]
        rows.append(cells)
    headers = rows[0] if rows else []
    data = [r for r in rows[1:] if not all(re.match(r"^[-:]+$", c) for c in r)]
    return {"headers": headers, "rows": data}


def parse_markdown(text: str) -> list[Slide]:
    body = strip_frontmatter(text)
    lines = body.split("\n")
    slides: list[Slide] = []
    current: Slide | None = None

    i = 0
    while i < len(lines):
        line = lines[i]

        if re.match(r"^---\s*$", line):
            if current:
                slides.append(current)
                current = None
            i += 1
            continue

        m = re.match(r"^#\s+(.+)$", line)
        if m:
            if current:
                slides.append(current)
            current = Slide(layout="title", title=m.group(1).strip())
            i += 1
            continue

        m = re.match(r"^##\s+(.+)$", line)
        if m:
            if current:
                slides.append(current)
            current = Slide(layout="content", title=m.group(1).strip())
            i += 1
            continue

        m = re.match(r"^###\s+(.+)$", line)
        if m:
            if not current:
                current = Slide(layout="content")
            current.blocks.append(("heading", m.group(1).strip()))
            i += 1
            continue

        if re.match(r"^\|.+\|", line):
            table_lines = []
            while i < len(lines) and re.match(r"^\|.+\|", lines[i]):
                table_lines.append(lines[i])
                i += 1
            if not current:
                current = Slide(layout="content")
            current.blocks.append(("table", _parse_table(table_lines)))
            continue

        m = re.match(r"^!\[([^\]]*)\]\(([^)]+)\)\s*$", line)
        if m:
            if not current:
                current = Slide(layout="content")
            current.blocks.append(("image", {"alt": m.group(1), "path": m.group(2)}))
            i += 1
            continue

        if re.match(r"^(\s*[-*+]|\s*\d+\.)\s", line):
            bullets: list[tuple[int, str]] = []
            while i < len(lines):
                bm = re.match(r"^(\s*)([-*+]|\d+\.)\s+(.+)$", lines[i])
                if bm:
                    level = min(len(bm.group(1)) // 2, 2)
                    bullets.append((level, bm.group(3).strip()))
                    i += 1
                elif lines[i].startswith("  ") and bullets:
                    bullets[-1] = (bullets[-1][0], bullets[-1][1] + " " + lines[i].strip())
                    i += 1
                else:
                    break
            if not current:
                current = Slide(layout="content")
            current.blocks.append(("bullets", bullets))
            continue

        if line.startswith("```"):
            code_lines: list[str] = []
            i += 1
            while i < len(lines) and not lines[i].startswith("```"):
                code_lines.append(lines[i])
                i += 1
            if i < len(lines):
                i += 1
            if not current:
                current = Slide(layout="content")
            current.blocks.append(("code", "\n".join(code_lines)))
            continue

        if line.strip():
            if not current:
                current = Slide(layout="content")
            current.blocks.append(("text", line.strip()))

        i += 1

    if current:
        slides.append(current)

    return slides


# ---------------------------------------------------------------------------
# PPTX renderer
# ---------------------------------------------------------------------------

def _find_layout(prs, name_hints: list[str], fallback_index: int):
    for layout in prs.slide_layouts:
        if layout.name.lower() in [h.lower() for h in name_hints]:
            return layout
    if fallback_index < len(prs.slide_layouts):
        return prs.slide_layouts[fallback_index]
    return prs.slide_layouts[0]


def _remove_template_slides(prs):
    from pptx.oxml.ns import qn

    for sldId in list(prs.slides._sldIdLst):
        rId = sldId.get(qn("r:id"))
        if rId:
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(sldId)


def _add_rich_runs(paragraph, text: str) -> None:
    parts = re.split(r"(\*\*[^*]+\*\*|\*[^*]+\*)", text)
    for part in parts:
        if not part:
            continue
        run = paragraph.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]
            run.font.bold = True
        elif part.startswith("*") and part.endswith("*"):
            run.text = part[1:-1]
            run.font.italic = True
        else:
            run.text = part


def _add_table(slide, table_data: dict) -> None:
    from pptx.util import Inches, Pt

    headers = table_data["headers"]
    rows = table_data["rows"]
    if not headers:
        return

    n_cols = len(headers)
    n_rows = len(rows) + 1

    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    row_height = Inches(0.35)

    table = slide.shapes.add_table(
        n_rows, n_cols, left, top, width, row_height * n_rows
    ).table

    col_width = int(Inches(9) / n_cols)
    for ci in range(n_cols):
        table.columns[ci].width = col_width

    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(11)

    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            if ci < n_cols:
                cell = table.cell(ri + 1, ci)
                cell.text = cell_text
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(10)


def render_pptx(
    slides: list[Slide],
    template_path: str | None,
    output_path: str,
    source_dir: str = ".",
) -> None:
    from pptx import Presentation
    from pptx.util import Emu, Inches, Pt

    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
        _remove_template_slides(prs)
    else:
        prs = Presentation()

    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(5143500)

    title_layout = _find_layout(prs, ["Title Slide", "title"], 0)
    content_layout = _find_layout(prs, ["Title and Content", "Title, Content"], 1)

    for slide_data in slides:
        if slide_data.layout == "title":
            slide = prs.slides.add_slide(title_layout)
            if slide.shapes.title:
                slide.shapes.title.text = slide_data.title
            subtitle_parts = [
                d for t, d in slide_data.blocks if t == "text"
            ]
            if subtitle_parts:
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 1:
                        ph.text = "\n".join(subtitle_parts)
                        break
        else:
            slide = prs.slides.add_slide(content_layout)
            if slide.shapes.title and slide_data.title:
                slide.shapes.title.text = slide_data.title

            body_ph = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    body_ph = ph
                    break

            if body_ph is not None:
                tf = body_ph.text_frame
                tf.clear()
                first = True

                for btype, bdata in slide_data.blocks:
                    if btype == "heading":
                        p = tf.paragraphs[0] if first else tf.add_paragraph()
                        first = False
                        run = p.add_run()
                        run.text = bdata
                        run.font.bold = True
                        run.font.size = Pt(18)

                    elif btype == "bullets":
                        for level, text in bdata:
                            p = tf.paragraphs[0] if first else tf.add_paragraph()
                            first = False
                            p.level = level
                            _add_rich_runs(p, text)

                    elif btype == "text":
                        p = tf.paragraphs[0] if first else tf.add_paragraph()
                        first = False
                        _add_rich_runs(p, bdata)

                    elif btype == "code":
                        p = tf.paragraphs[0] if first else tf.add_paragraph()
                        first = False
                        run = p.add_run()
                        run.text = bdata
                        run.font.name = "Courier New"
                        run.font.size = Pt(10)

                    elif btype == "image":
                        img_path = Path(source_dir) / bdata["path"]
                        if img_path.exists():
                            slide.shapes.add_picture(
                                str(img_path),
                                Inches(1),
                                Inches(3),
                                width=Inches(8),
                            )

                    elif btype == "table":
                        _add_table(slide, bdata)

    prs.save(output_path)


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------


def main() -> None:
    parser = argparse.ArgumentParser(description="Render markdown to PPTX")
    parser.add_argument("input", help="Input markdown file")
    parser.add_argument("output", help="Output PPTX file")
    parser.add_argument("--template", help="PPTX template file")
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Input not found: {args.input}", file=sys.stderr)
        sys.exit(2)

    text = input_path.read_text(errors="replace")
    slides = parse_markdown(text)

    if not slides:
        slides = [Slide(layout="title", title="(empty presentation)")]

    render_pptx(slides, args.template, args.output, str(input_path.parent))


if __name__ == "__main__":
    main()
