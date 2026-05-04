#!/usr/bin/env python3
"""
Generate brand/templates/ from brand/tokens.json.

Creates:
  brand/templates/reference.docx   DOCX reference-doc with brand heading colours
  brand/templates/template.pptx    PowerPoint with brand theme + title slide
  brand/templates/template.xlsx    Excel with brand-coloured header row

Requires: python-docx, python-pptx, openpyxl
Run scripts/install-local.sh to install all three.

tokens.json `docx` block (all optional — defaults are applied when absent):

  {
    "docx": {
      "margins_cm": 2.5,
      "body": {
        "color": "#2A3340",
        "size": 11,
        "line_spacing": 1.25,
        "space_after": 8
      },
      "title":     { "size": 32 },
      "heading_1": { "size": 22, "before": 18, "after": 8 },
      "heading_2": { "size": 16, "before": 14, "after": 6 },
      "heading_3": { "size": 13, "before": 10, "after": 4 },
      "heading_4": { "size": 11, "before": 8,  "after": 3 },
      "footer":    "{name}  •  Confidentiel"
    }
  }
"""

from __future__ import annotations

import json
import subprocess
import sys
import tempfile
from pathlib import Path

ROOT = Path.cwd()
TEMPLATES = ROOT / "brand" / "templates"


# ---------------------------------------------------------------------------
# Colour helpers
# ---------------------------------------------------------------------------

def _hex_to_rgb(h: str) -> tuple[int, int, int]:
    h = h.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _darken(h: str, pct: float = 0.25) -> str:
    r, g, b = _hex_to_rgb(h)
    return "#{:02x}{:02x}{:02x}".format(
        int(r * (1 - pct)), int(g * (1 - pct)), int(b * (1 - pct))
    )


def _lighten(h: str, pct: float = 0.85) -> str:
    r, g, b = _hex_to_rgb(h)
    return "#{:02x}{:02x}{:02x}".format(
        int(r + (255 - r) * pct),
        int(g + (255 - g) * pct),
        int(b + (255 - b) * pct),
    )


# ---------------------------------------------------------------------------
# DOCX default constants (BPI/investor-grade)
# ---------------------------------------------------------------------------

_DOCX_DEFAULTS = {
    "margins_cm": 2.5,
    "body": {
        "color": "#2A3340",
        "size": 11,
        "line_spacing": 1.25,
        "space_after": 8,
    },
    "title":     {"size": 32},
    "heading_1": {"size": 22, "before": 18, "after": 8},
    "heading_2": {"size": 16, "before": 14, "after": 6},
    "heading_3": {"size": 13, "before": 10, "after": 4},
    "heading_4": {"size": 11, "before": 8,  "after": 3},
    "footer":    "{name}  •  Confidentiel",
}


def _merge_docx_cfg(tokens: dict) -> dict:
    """Deep-merge tokens['docx'] over _DOCX_DEFAULTS."""
    import copy
    cfg = copy.deepcopy(_DOCX_DEFAULTS)
    user = tokens.get("docx", {})
    for key, val in user.items():
        if isinstance(val, dict) and isinstance(cfg.get(key), dict):
            cfg[key].update(val)
        else:
            cfg[key] = val
    return cfg


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def _generate_docx(tokens: dict) -> None:
    from docx import Document  # type: ignore
    from docx.shared import Cm, Pt, RGBColor  # type: ignore
    from docx.enum.text import WD_LINE_SPACING  # type: ignore

    primary = tokens["colors"]["primary"]
    font_name = tokens["fonts"]["primary"]
    project_name = tokens.get("name", "Project")

    cfg = _merge_docx_cfg(tokens)

    r, g, b = _hex_to_rgb(primary)
    h2_r, h2_g, h2_b = _hex_to_rgb(_darken(primary, 0.15))
    h3_r, h3_g, h3_b = _hex_to_rgb(_darken(primary, 0.35))
    # Muted slate for H4 — fixed colour regardless of brand primary
    h4_r, h4_g, h4_b = 0x6B, 0x74, 0x80  # slate-500 equivalent

    # Bootstrap from pandoc's default reference.docx so the XML structure
    # is already what pandoc expects.
    doc: Document
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
        tmp_path = Path(tmp.name)

    result = subprocess.run(
        ["pandoc", "--print-default-data-file", "reference.docx"],
        capture_output=True,
    )
    if result.returncode == 0 and result.stdout:
        tmp_path.write_bytes(result.stdout)
        doc = Document(str(tmp_path))
    else:
        doc = Document()
    tmp_path.unlink(missing_ok=True)

    # --- Page margins ---
    margins_cm = cfg["margins_cm"]
    for section in doc.sections:
        section.left_margin   = Cm(margins_cm)
        section.right_margin  = Cm(margins_cm)
        section.top_margin    = Cm(margins_cm)
        section.bottom_margin = Cm(margins_cm)

    # --- Normal body style ---
    body_cfg = cfg["body"]
    body_color = body_cfg.get("color", "#2A3340")
    bc_r, bc_g, bc_b = _hex_to_rgb(body_color)
    try:
        normal = doc.styles["Normal"]
        normal.font.size = Pt(body_cfg.get("size", 11))
        normal.font.name = font_name
        normal.font.color.rgb = RGBColor(bc_r, bc_g, bc_b)
        normal.paragraph_format.space_after = Pt(body_cfg.get("space_after", 8))
        normal.paragraph_format.line_spacing = body_cfg.get("line_spacing", 1.25)
    except (KeyError, Exception):
        pass

    # --- Title style ---
    title_cfg = cfg.get("title", {})
    try:
        title_style = doc.styles["Title"]
        title_style.font.size = Pt(title_cfg.get("size", 32))
        title_style.font.name = font_name
        title_style.paragraph_format.space_after = Pt(18)
    except (KeyError, Exception):
        pass

    # --- Heading styles (H1–H4) ---
    heading_defs = [
        ("Heading 1", (r, g, b),             cfg["heading_1"], True),
        ("Heading 2", (h2_r, h2_g, h2_b),    cfg["heading_2"], True),
        ("Heading 3", (h3_r, h3_g, h3_b),    cfg["heading_3"], True),
        ("Heading 4", (h4_r, h4_g, h4_b),    cfg["heading_4"], False),
    ]
    for style_name, rgb, hcfg, bold in heading_defs:
        try:
            style = doc.styles[style_name]
            style.font.color.rgb = RGBColor(*rgb)
            style.font.bold = bold
            style.font.size = Pt(hcfg.get("size", 11))
            style.font.name = font_name
            if style_name == "Heading 4":
                style.font.italic = True
            pf = style.paragraph_format
            before = hcfg.get("before")
            after  = hcfg.get("after")
            if before is not None:
                pf.space_before = Pt(before)
            if after is not None:
                pf.space_after = Pt(after)
        except (KeyError, Exception):
            pass

    # --- Footer ---
    footer_tpl = cfg.get("footer", "{name}  •  Confidentiel")
    footer_text = footer_tpl.replace("{name}", project_name)
    try:
        for section in doc.sections:
            footer = section.footer
            # Clear existing paragraphs content
            for para in footer.paragraphs:
                para.clear()
                run = para.add_run(footer_text)
                run.font.size = Pt(8)
                run.font.name = font_name
                from docx.enum.text import WD_ALIGN_PARAGRAPH  # type: ignore
                para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                break  # only set the first paragraph
    except (KeyError, Exception):
        pass

    out = TEMPLATES / "reference.docx"
    doc.save(str(out))
    print(f"  ✓ {out.relative_to(ROOT)}")


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def _generate_pptx(tokens: dict) -> None:
    from pptx import Presentation  # type: ignore
    from pptx.dml.color import RGBColor  # type: ignore
    from pptx.oxml.ns import qn  # type: ignore
    from pptx.util import Emu, Pt  # type: ignore
    from lxml import etree  # type: ignore

    primary = tokens["colors"]["primary"]
    name = tokens.get("name", "Project")
    primary_no_hash = primary.lstrip("#")

    prs = Presentation()
    prs.slide_width  = Emu(9144000)   # 10 in
    prs.slide_height = Emu(5143500)   # 5.625 in

    # Apply primary to the slide master's theme colour scheme (accent1 + dk1).
    for master in prs.slide_masters:
        clr_scheme = master.element.find(".//" + qn("a:clrScheme"))
        if clr_scheme is None:
            continue
        for tag in (qn("a:dk1"), qn("a:accent1")):
            elem = clr_scheme.find(tag)
            if elem is None:
                continue
            # Remove existing colour children (sysClr / srgbClr) and replace.
            for child in list(elem):
                elem.remove(child)
            srgb = etree.SubElement(elem, qn("a:srgbClr"))
            srgb.set("val", primary_no_hash)

    # Title slide — layout 0 is always "Title Slide" in the default theme.
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_ph = slide.shapes.title
    if title_ph is not None:
        title_ph.text = name
        run = title_ph.text_frame.paragraphs[0].runs[0]
        run.font.color.rgb = RGBColor(*_hex_to_rgb(primary))
        run.font.size = Pt(36)

    out = TEMPLATES / "template.pptx"
    prs.save(str(out))
    print(f"  ✓ {out.relative_to(ROOT)}")


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------

def _generate_xlsx(tokens: dict) -> None:
    from openpyxl import Workbook  # type: ignore
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side  # type: ignore

    primary = tokens["colors"]["primary"].lstrip("#")
    font_name = tokens["fonts"]["primary"]
    light_hex = _lighten("#" + primary, 0.93).lstrip("#")

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"

    header_fill   = PatternFill("solid", fgColor=primary)
    header_font   = Font(bold=True, color="FFFFFF", name=font_name)
    body_font     = Font(name=font_name)
    alt_fill      = PatternFill("solid", fgColor=light_hex)
    thin          = Side(style="thin", color="DDDDDD")
    border        = Border(left=thin, right=thin, top=thin, bottom=thin)
    center        = Alignment(horizontal="center", vertical="center")

    for col, header in enumerate(["Column A", "Column B", "Column C"], start=1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.fill   = header_fill
        cell.font   = header_font
        cell.border = border
        cell.alignment = center
        ws.column_dimensions[cell.column_letter].width = 20

    ws.row_dimensions[1].height = 22

    for row in range(2, 6):
        for col in range(1, 4):
            cell = ws.cell(row=row, column=col, value=f"Sample {row - 1}-{col}")
            cell.font   = body_font
            cell.border = border
            if row % 2 == 0:
                cell.fill = alt_fill

    out = TEMPLATES / "template.xlsx"
    wb.save(str(out))
    print(f"  ✓ {out.relative_to(ROOT)}")


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main() -> None:
    tokens_path = ROOT / "brand" / "tokens.json"
    if not tokens_path.exists():
        print(f"Error: {tokens_path} not found. Run scan-brand.py first.", file=sys.stderr)
        sys.exit(1)

    tokens = json.loads(tokens_path.read_text())
    TEMPLATES.mkdir(parents=True, exist_ok=True)

    skipped: list[str] = []
    failed: list[str] = []
    succeeded: list[str] = []

    for label, fn in [
        ("docx",  _generate_docx),
        ("pptx",  _generate_pptx),
        ("xlsx",  _generate_xlsx),
    ]:
        try:
            fn(tokens)
            succeeded.append(label)
        except ImportError as e:
            skipped.append(label)
            print(f"  ⚠️  {label.upper()}: missing dep — {e}. Run scripts/install-local.sh.", file=sys.stderr)
        except Exception as e:
            failed.append(label)
            print(f"  ⚠️  {label.upper()}: generation failed — {e}", file=sys.stderr)

    parts = (
        [f"{f} ok" for f in succeeded]
        + [f"{f} skipped (missing dep)" for f in skipped]
        + [f"{f} failed" for f in failed]
    )
    print("templates: " + ", ".join(parts))

    if not succeeded:
        sys.exit(1)


if __name__ == "__main__":
    main()
