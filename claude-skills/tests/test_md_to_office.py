#!/usr/bin/env python3
"""
Unit tests for the md-to-office skill.

Covers four layers, from cheap to expensive:

  1. Pure shell logic of resolve-template.sh — five-priority chain from
     PRD-008 §5.3. Fast; runs without external deps.
  2. Orchestrator semantics of md-to-office.sh — target gating, output
     path convention, idempotency, --force bypass. Exercises bash +
     the skill's own scripts; also runs without external deps because
     the docx renderer is stubbed out via PATH shim.
  3. Brand scanner (scan-brand.py) — pure Python stdlib; tests each
     signal source (README, CSS, tailwind, package.json, etc.) and
     merge / fallback behaviour.
  4. End-to-end tests — pandoc / python-docx / python-pptx / openpyxl
     invoked for real. Skipped when deps are missing so CI doesn't
     break for devs who haven't run install-local.sh yet.

Run locally:

    python3 claude-skills/tests/test_md_to_office.py

Stdlib only for layers 1-3. Layer 4 auto-skips when deps are absent.
"""

from __future__ import annotations

import json
import os
import shutil
import subprocess
import sys
import tempfile
import textwrap
import time
import unittest
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[2]
SKILL_DIR = REPO_ROOT / "claude-skills" / "skills" / "md-to-office"
SCRIPTS = SKILL_DIR / "scripts"

RESOLVE       = SCRIPTS / "resolve-template.sh"
ORCH          = SCRIPTS / "md-to-office.sh"
RENDER_DOCX   = SCRIPTS / "render-docx.sh"
RENDER_PPTX   = SCRIPTS / "render-pptx.sh"
RENDER_PPTX_PY = SCRIPTS / "render-pptx.py"
SCAN_BRAND    = SCRIPTS / "scan-brand.py"
GEN_TEMPLATES = SCRIPTS / "generate-templates.py"
INIT_BRAND    = SCRIPTS / "init-brand.sh"

# Magic bytes of a ZIP archive — DOCX/PPTX/XLSX all use the ZIP container.
ZIP_MAGIC = b"PK\x03\x04"


def run(
    cmd: list[str],
    *,
    cwd: Path | None = None,
    env: dict[str, str] | None = None,
    check: bool = True,
) -> subprocess.CompletedProcess[str]:
    """Run a command, capturing stdout/stderr as text."""
    final_env = os.environ.copy()
    if env:
        final_env.update(env)
    return subprocess.run(
        cmd,
        cwd=str(cwd) if cwd else None,
        env=final_env,
        text=True,
        capture_output=True,
        check=check,
    )


class TestResolveTemplate(unittest.TestCase):
    """Covers the 5-priority resolution chain from PRD-008 §5.3."""

    def setUp(self) -> None:
        self.tmp = Path(tempfile.mkdtemp(prefix="bsg-mdto-"))
        (self.tmp / "brand" / "templates").mkdir(parents=True)

    def tearDown(self) -> None:
        shutil.rmtree(self.tmp, ignore_errors=True)

    def _touch(self, rel: str) -> Path:
        p = self.tmp / rel
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_bytes(ZIP_MAGIC + b"stub")
        return p

    def _resolve(
        self, target: str, *extra: str, env: dict[str, str] | None = None
    ) -> str:
        result = run(
            ["bash", str(RESOLVE), target, *extra],
            cwd=self.tmp,
            env=env,
        )
        return result.stdout.strip()

    # ---- level 1 -----------------------------------------------------------

    def test_flag_override_wins(self) -> None:
        override = self._touch("custom/foo.docx")
        # Also drop a canonical template to prove the override is preferred.
        self._touch("brand/templates/docx.docx")
        out = self._resolve("docx", "--override", str(override))
        self.assertEqual(out, str(override))

    def test_flag_override_missing_is_error(self) -> None:
        result = run(
            ["bash", str(RESOLVE), "docx", "--override", "does/not/exist.docx"],
            cwd=self.tmp,
            check=False,
        )
        self.assertNotEqual(result.returncode, 0)
        self.assertIn("not found", result.stderr.lower())

    # ---- level 2 -----------------------------------------------------------

    def test_env_var_wins_over_canonical(self) -> None:
        env_dir = self.tmp / "env-templates"
        env_dir.mkdir()
        env_template = env_dir / "docx.docx"
        env_template.write_bytes(ZIP_MAGIC + b"envstub")
        # Canonical exists too — env var should outrank it.
        self._touch("brand/templates/docx.docx")
        out = self._resolve(
            "docx",
            env={"BSG_BRAND_TEMPLATES": str(env_dir)},
        )
        self.assertEqual(out, str(env_template))

    # ---- level 3 -----------------------------------------------------------

    def test_canonical_per_repo_template(self) -> None:
        tmpl = self._touch("brand/templates/docx.docx")
        out = self._resolve("docx")
        self.assertEqual(out, "./brand/templates/docx.docx")
        self.assertTrue(tmpl.exists())  # sanity

    # ---- level 4 -----------------------------------------------------------

    def test_legacy_reference_docx_alias(self) -> None:
        self._touch("brand/templates/reference.docx")
        out = self._resolve("docx")
        self.assertEqual(out, "./brand/templates/reference.docx")

    def test_legacy_alias_does_not_apply_to_pptx(self) -> None:
        # reference.docx exists but target is pptx — the alias must not leak.
        self._touch("brand/templates/reference.docx")
        out = self._resolve("pptx")
        self.assertEqual(out, "")

    # ---- level 5 -----------------------------------------------------------

    def test_no_template_returns_empty_success(self) -> None:
        """No brand/, no env, no override — exit 0 with empty stdout."""
        result = run(["bash", str(RESOLVE), "docx"], cwd=self.tmp)
        self.assertEqual(result.returncode, 0)
        self.assertEqual(result.stdout.strip(), "")

    # ---- argument validation -----------------------------------------------

    def test_unknown_target_is_error(self) -> None:
        result = run(
            ["bash", str(RESOLVE), "html"],
            cwd=self.tmp,
            check=False,
        )
        self.assertNotEqual(result.returncode, 0)

    def test_missing_target_is_error(self) -> None:
        result = run(
            ["bash", str(RESOLVE)],
            cwd=self.tmp,
            check=False,
        )
        self.assertNotEqual(result.returncode, 0)


class TestOrchestrator(unittest.TestCase):
    """Covers md-to-office.sh behaviors that don't need pandoc.

    A fake `pandoc` is placed on PATH and just emits a stub DOCX so we
    can verify the orchestrator calls the renderer with the right args,
    resolves templates correctly, and honors idempotency / --force.
    """

    def setUp(self) -> None:
        self.tmp = Path(tempfile.mkdtemp(prefix="bsg-orch-"))
        self.bin = self.tmp / "bin"
        self.bin.mkdir()
        fake = self.bin / "pandoc"
        fake.write_text(
            textwrap.dedent(
                """\
                #!/usr/bin/env bash
                # Find the -o <output> pair and emit a stub ZIP there.
                out=""
                while [[ $# -gt 0 ]]; do
                  case "$1" in
                    -o) out="$2"; shift 2 ;;
                    *)  shift ;;
                  esac
                done
                printf 'PK\\x03\\x04stub-from-fake-pandoc' > "$out"
                """
            )
        )
        fake.chmod(0o755)
        self.env = {"PATH": f"{self.bin}:{os.environ.get('PATH', '')}"}

        self.src = self.tmp / "report.md"
        self.src.write_text("# Hello\n\nBody.\n")

    def tearDown(self) -> None:
        shutil.rmtree(self.tmp, ignore_errors=True)

    def _run_orch(self, *args: str, check: bool = True) -> subprocess.CompletedProcess[str]:
        return run(
            ["bash", str(ORCH), *args],
            cwd=self.tmp,
            env=self.env,
            check=check,
        )

    # ---- output path convention --------------------------------------------

    def test_default_output_path_swaps_extension(self) -> None:
        result = self._run_orch(str(self.src))
        out = self.tmp / "report.docx"
        self.assertTrue(out.exists(), f"expected {out}; stderr={result.stderr}")
        self.assertTrue(out.read_bytes().startswith(ZIP_MAGIC))
        self.assertIn(str(out), result.stdout)

    def test_explicit_out_flag_is_honored(self) -> None:
        custom = self.tmp / "dist" / "custom.docx"
        self._run_orch(str(self.src), "--out", str(custom))
        self.assertTrue(custom.exists())

    # ---- idempotency -------------------------------------------------------

    def test_idempotent_when_output_is_newer(self) -> None:
        out = self.tmp / "report.docx"
        out.write_bytes(ZIP_MAGIC + b"existing")
        # Make the output newer than the source.
        future = time.time() + 10
        os.utime(out, (future, future))

        result = self._run_orch(str(self.src))
        self.assertIn("cached", result.stdout)
        # Body unchanged — fake pandoc would have overwritten it.
        self.assertEqual(out.read_bytes(), ZIP_MAGIC + b"existing")

    def test_force_flag_bypasses_cache(self) -> None:
        out = self.tmp / "report.docx"
        out.write_bytes(ZIP_MAGIC + b"existing")
        future = time.time() + 10
        os.utime(out, (future, future))

        self._run_orch(str(self.src), "--force")
        # Fake pandoc overwrites with its stub payload.
        self.assertIn(b"stub-from-fake-pandoc", out.read_bytes())

    # ---- target gating -----------------------------------------------------

    def test_pptx_target_accepted(self) -> None:
        """pptx is now a valid target; the orchestrator should not reject it."""
        try:
            __import__("pptx")
        except ImportError:
            self.skipTest("python-pptx not installed")
        result = self._run_orch(str(self.src), "--target", "pptx", check=False)
        self.assertEqual(result.returncode, 0)
        out = self.tmp / "report.pptx"
        self.assertTrue(out.exists(), f"expected {out}; stderr={result.stderr}")

    def test_xlsx_target_is_not_implemented(self) -> None:
        result = self._run_orch(str(self.src), "--target", "xlsx", check=False)
        self.assertNotEqual(result.returncode, 0)

    def test_unknown_target_is_error(self) -> None:
        result = self._run_orch(str(self.src), "--target", "html", check=False)
        self.assertNotEqual(result.returncode, 0)

    # ---- auto-detection from frontmatter -----------------------------------

    def test_autodetect_gamma_presentation_produces_pptx(self) -> None:
        try:
            __import__("pptx")
        except ImportError:
            self.skipTest("python-pptx not installed")
        src = self.tmp / "deck.md"
        src.write_text("---\ngamma:\n  format: presentation\n---\n# Deck\n\nSlide content.\n")
        result = self._run_orch(str(src), check=False)
        self.assertEqual(result.returncode, 0, f"stderr={result.stderr}")
        self.assertTrue((self.tmp / "deck.pptx").exists())

    def test_autodetect_gamma_document_produces_docx(self) -> None:
        src = self.tmp / "report.md"
        src.write_text("---\ngamma:\n  format: document\n---\n# Report\n\nBody.\n")
        result = self._run_orch(str(src))
        self.assertTrue((self.tmp / "report.docx").exists())

    def test_autodetect_target_field_in_frontmatter(self) -> None:
        try:
            __import__("pptx")
        except ImportError:
            self.skipTest("python-pptx not installed")
        src = self.tmp / "slides.md"
        src.write_text("---\ntarget: pptx\n---\n# Slides\n\nContent.\n")
        result = self._run_orch(str(src), check=False)
        self.assertEqual(result.returncode, 0, f"stderr={result.stderr}")
        self.assertTrue((self.tmp / "slides.pptx").exists())

    def test_autodetect_no_frontmatter_defaults_to_docx(self) -> None:
        result = self._run_orch(str(self.src))
        self.assertTrue((self.tmp / "report.docx").exists())

    def test_explicit_target_overrides_frontmatter(self) -> None:
        src = self.tmp / "deck.md"
        src.write_text("---\ngamma:\n  format: presentation\n---\n# Deck\n")
        result = self._run_orch(str(src), "--target", "docx")
        self.assertTrue((self.tmp / "deck.docx").exists())

    # ---- batch mode --------------------------------------------------------

    def test_batch_converts_directory(self) -> None:
        batch_dir = self.tmp / "docs"
        batch_dir.mkdir()
        (batch_dir / "a.md").write_text("# Doc A\n\nBody.\n")
        (batch_dir / "b.md").write_text("# Doc B\n\nBody.\n")
        result = self._run_orch(str(batch_dir))
        self.assertTrue((batch_dir / "a.docx").exists())
        self.assertTrue((batch_dir / "b.docx").exists())

    def test_batch_autodetects_per_file(self) -> None:
        try:
            __import__("pptx")
        except ImportError:
            self.skipTest("python-pptx not installed")
        batch_dir = self.tmp / "mixed"
        batch_dir.mkdir()
        (batch_dir / "deck.md").write_text("---\ngamma:\n  format: presentation\n---\n# Deck\n")
        (batch_dir / "doc.md").write_text("---\ngamma:\n  format: document\n---\n# Doc\n")
        result = self._run_orch(str(batch_dir), check=False)
        self.assertEqual(result.returncode, 0, f"stderr={result.stderr}")
        self.assertTrue((batch_dir / "deck.pptx").exists())
        self.assertTrue((batch_dir / "doc.docx").exists())

    def test_batch_empty_dir_is_error(self) -> None:
        empty = self.tmp / "empty"
        empty.mkdir()
        result = self._run_orch(str(empty), check=False)
        self.assertNotEqual(result.returncode, 0)

    # ---- missing input -----------------------------------------------------

    def test_missing_input_is_error(self) -> None:
        result = self._run_orch(str(self.tmp / "nope.md"), check=False)
        self.assertNotEqual(result.returncode, 0)


class TestScanBrand(unittest.TestCase):
    """Covers scan-brand.py signal sources. Pure Python stdlib, no external deps."""

    def setUp(self) -> None:
        self.tmp = Path(tempfile.mkdtemp(prefix="bsg-scan-"))

    def tearDown(self) -> None:
        shutil.rmtree(self.tmp, ignore_errors=True)

    def _scan(self, env: dict[str, str] | None = None) -> dict:
        result = run(
            ["python3", str(SCAN_BRAND)],
            cwd=self.tmp,
            env={**os.environ, **(env or {})},
        )
        return json.loads(result.stdout)

    # ---- name detection ----------------------------------------------------

    def test_name_from_readme_h1(self) -> None:
        (self.tmp / "README.md").write_text("# Acme Platform\n\nDesc.\n")
        self.assertEqual(self._scan()["name"], "Acme Platform")

    def test_name_from_package_json(self) -> None:
        (self.tmp / "package.json").write_text('{"name": "acme-frontend", "version": "1.0.0"}')
        self.assertEqual(self._scan()["name"], "acme-frontend")

    def test_name_from_pyproject_toml(self) -> None:
        (self.tmp / "pyproject.toml").write_text('[project]\nname = "acme-backend"\n')
        self.assertEqual(self._scan()["name"], "acme-backend")

    def test_name_from_stack_yaml(self) -> None:
        (self.tmp / "stack.yaml").write_text("name: acme-stack\nversion: 1\n")
        self.assertEqual(self._scan()["name"], "acme-stack")

    def test_name_from_narrative_wins_over_readme(self) -> None:
        (self.tmp / "README.md").write_text("# Readme Name\n")
        brand = self.tmp / "brand"
        brand.mkdir()
        (brand / "NARRATIVE.md").write_text("# Narrative Name\n\nBrand voice.\n")
        self.assertEqual(self._scan()["name"], "Narrative Name")

    def test_name_defaults_to_project(self) -> None:
        # Totally empty repo with no git — should not crash.
        tokens = self._scan()
        self.assertIn("name", tokens)
        self.assertTrue(len(tokens["name"]) > 0)

    # ---- primary colour detection ------------------------------------------

    def test_primary_from_css_custom_property(self) -> None:
        css = self.tmp / "styles.css"
        css.write_text(":root { --primary-color: #e63946; }\n")
        self.assertEqual(self._scan()["colors"]["primary"], "#e63946")

    def test_primary_from_css_color_primary(self) -> None:
        css = self.tmp / "theme.css"
        css.write_text(":root { --color-primary: #2b9348; }\n")
        self.assertEqual(self._scan()["colors"]["primary"], "#2b9348")

    def test_primary_from_tailwind_config(self) -> None:
        tw = self.tmp / "tailwind.config.js"
        tw.write_text("module.exports = { theme: { colors: { primary: '#6d28d9' } } }\n")
        self.assertEqual(self._scan()["colors"]["primary"], "#6d28d9")

    def test_primary_from_tailwind_v4_theme_block(self) -> None:
        """#241: Tailwind v4 @theme blocks with --color-* tokens must be parsed."""
        apps_css = self.tmp / "apps" / "web-app" / "src" / "app"
        apps_css.mkdir(parents=True)
        (apps_css / "globals.css").write_text(
            "@theme {\n"
            "  --color-ef-blue-50: #eff6ff;\n"
            "  --color-ef-blue-500: #1A56DB;\n"
            "  --color-ef-blue-900: #1e3a5f;\n"
            "}\n"
        )
        result = self._scan()["colors"]["primary"]
        self.assertEqual(
            result.lower(),
            "#1a56db",
            f"Expected #1a56db from Tailwind v4 @theme block, got {result}",
        )

    def test_primary_app_css_preferred_over_public_html(self) -> None:
        """#241: main app CSS should win over stray HTML files under public/."""
        # Stray demo HTML in public/ with a different color
        pub = self.tmp / "public" / "images" / "brand"
        pub.mkdir(parents=True)
        (pub / "index.html").write_text(
            "<link href='https://fonts.googleapis.com/css2?family=Outfit' rel='stylesheet'>\n"
            "<style>:root { --primary-color: #badcol; }</style>\n"
        )
        # Real app globals.css with the actual brand color
        apps_css = self.tmp / "apps" / "web-app" / "src" / "app"
        apps_css.mkdir(parents=True)
        (apps_css / "globals.css").write_text(
            "@theme {\n"
            "  --color-ef-blue-500: #1A56DB;\n"
            "}\n"
        )
        result = self._scan()["colors"]["primary"]
        self.assertEqual(
            result.lower(),
            "#1a56db",
            f"Expected #1a56db from app globals.css, got {result} (stray public HTML color leaked)",
        )


    def test_primary_from_tokens_json(self) -> None:
        (self.tmp / "tokens.json").write_text(
            json.dumps({"colors": {"primary": "#ff6b35"}})
        )
        self.assertEqual(self._scan()["colors"]["primary"], "#ff6b35")

    def test_primary_shorthand_hex_expanded(self) -> None:
        css = self.tmp / "main.css"
        css.write_text(":root { --primary: #f0f; }\n")
        result = self._scan()["colors"]["primary"]
        self.assertEqual(result, "#ff00ff")

    def test_primary_defaults_to_bsg_navy(self) -> None:
        tokens = self._scan()
        self.assertEqual(tokens["colors"]["primary"], "#1a1a2e")

    # ---- font detection ----------------------------------------------------

    def test_font_from_css_variable(self) -> None:
        css = self.tmp / "vars.css"
        css.write_text(":root { --font-family-primary: 'Inter', sans-serif; }\n")
        font = self._scan()["fonts"]["primary"]
        self.assertIn("Inter", font)

    def test_font_defaults_to_helvetica(self) -> None:
        self.assertEqual(self._scan()["fonts"]["primary"], "Helvetica Neue")

    # ---- merge with existing partial tokens --------------------------------

    def test_existing_tokens_take_priority_over_scan(self) -> None:
        brand = self.tmp / "brand"
        brand.mkdir()
        existing = {"name": "Locked Name", "colors": {"primary": "#locked"}, "fonts": {"primary": "Locked Font"}}
        (brand / "tokens.json").write_text(json.dumps(existing))
        # Also put a README that would normally win.
        (self.tmp / "README.md").write_text("# Different Name\n")
        tokens = self._scan()
        self.assertEqual(tokens["name"], "Locked Name")
        self.assertEqual(tokens["colors"]["primary"], "#locked")

    def test_partial_tokens_filled_by_scan(self) -> None:
        brand = self.tmp / "brand"
        brand.mkdir()
        (brand / "tokens.json").write_text(json.dumps({"name": "My Project"}))
        (self.tmp / "styles.css").write_text(":root { --primary: #aabbcc; }\n")
        tokens = self._scan()
        self.assertEqual(tokens["name"], "My Project")   # from existing
        self.assertEqual(tokens["colors"]["primary"], "#aabbcc")  # from CSS

    # ---- @font-face detection ------------------------------------------------

    def test_font_from_font_face(self) -> None:
        css = self.tmp / "fonts.css"
        css.write_text(
            "@font-face {\n"
            "  font-family: 'Cabinet Grotesk';\n"
            "  src: url('cabinet.woff2');\n"
            "}\n"
        )
        self.assertEqual(self._scan()["fonts"]["primary"], "Cabinet Grotesk")

    def test_font_from_google_fonts_import(self) -> None:
        css = self.tmp / "global.css"
        css.write_text(
            "@import url('https://fonts.googleapis.com/css2?family=Poppins:wght@400;700');\n"
        )
        self.assertEqual(self._scan()["fonts"]["primary"], "Poppins")

    def test_css_variable_font_wins_over_font_face(self) -> None:
        css = self.tmp / "all.css"
        css.write_text(
            ":root { --font-family-primary: 'Manrope', sans-serif; }\n"
            "@font-face { font-family: 'Cabinet Grotesk'; src: url('c.woff2'); }\n"
        )
        self.assertIn("Manrope", self._scan()["fonts"]["primary"])

    # ---- logo detection ------------------------------------------------------

    def test_logos_found_in_public_dir(self) -> None:
        pub = self.tmp / "public"
        pub.mkdir()
        (pub / "logo.svg").write_text("<svg/>")
        (pub / "logo-dark.png").write_bytes(b"\x89PNG")
        tokens = self._scan()
        self.assertIn("logos", tokens)
        self.assertEqual(len(tokens["logos"]), 2)
        self.assertTrue(any("logo.svg" in l for l in tokens["logos"]))

    def test_logos_empty_when_none_found(self) -> None:
        tokens = self._scan()
        self.assertEqual(tokens["logos"], [])

    # ---- identity doc detection ----------------------------------------------

    def test_identity_docs_found(self) -> None:
        brand = self.tmp / "brand"
        brand.mkdir()
        (brand / "NARRATIVE.md").write_text("# Brand\n\nVoice.\n")
        (self.tmp / "DESIGN.md").write_text("# Design System\n")
        tokens = self._scan()
        self.assertIn("identity_docs", tokens)
        self.assertTrue(len(tokens["identity_docs"]) >= 2)

    def test_identity_docs_empty_when_none(self) -> None:
        tokens = self._scan()
        self.assertEqual(tokens["identity_docs"], [])

    # ---- existing template detection -----------------------------------------

    def test_existing_templates_found(self) -> None:
        (self.tmp / "old-report.docx").write_bytes(b"PK\x03\x04stub")
        tokens = self._scan()
        self.assertIn("existing_templates", tokens)
        self.assertTrue(any("old-report.docx" in t for t in tokens["existing_templates"]))

    # ---- audit generation ----------------------------------------------------

    def test_audit_flag_writes_markdown(self) -> None:
        (self.tmp / "README.md").write_text("# Audit Test\n\nDesc.\n")
        css = self.tmp / "style.css"
        css.write_text(":root { --primary-color: #abcdef; }\n")
        audit_path = self.tmp / "brand" / "templates" / "brand-audit.md"
        result = run(
            ["python3", str(SCAN_BRAND), "--audit", str(audit_path)],
            cwd=self.tmp,
        )
        self.assertTrue(audit_path.exists(), f"audit not written; stderr={result.stderr}")
        content = audit_path.read_text()
        self.assertIn("# Brand Audit", content)
        self.assertIn("Audit Test", content)
        self.assertIn("#abcdef", content)
        self.assertIn("Chosen tokens", content)
        self.assertIn("Discovery details", content)

    def test_audit_shows_defaults_when_empty_repo(self) -> None:
        audit_path = self.tmp / "audit.md"
        run(
            ["python3", str(SCAN_BRAND), "--audit", str(audit_path)],
            cwd=self.tmp,
        )
        content = audit_path.read_text()
        self.assertIn("Defaults applied", content)

    # ---- orchestrator onboarding banner ------------------------------------

    def _make_fake_pandoc(self, bin_dir: Path) -> None:
        """Write a portable fake pandoc that creates a stub ZIP at the -o path."""
        fake = bin_dir / "pandoc"
        fake.write_text(textwrap.dedent("""\
            #!/usr/bin/env bash
            while [[ $# -gt 0 ]]; do
              case "$1" in
                -o) printf 'PK\\x03\\x04stub' > "$2"; shift 2 ;;
                *)  shift ;;
              esac
            done
            """))
        fake.chmod(0o755)

    def test_orchestrator_shows_init_hint_when_no_brand_dir(self) -> None:
        src = self.tmp / "doc.md"
        src.write_text("# Hello\n\nBody.\n")
        bin_dir = self.tmp / "bin"
        bin_dir.mkdir()
        self._make_fake_pandoc(bin_dir)
        result = run(
            ["bash", str(ORCH), str(src)],
            cwd=self.tmp,
            env={**os.environ, "PATH": f"{bin_dir}:{os.environ.get('PATH', '')}"},
        )
        self.assertIn("--init", result.stderr)
        self.assertIn("No brand standard", result.stderr)

    def test_orchestrator_renders_despite_missing_brand(self) -> None:
        """Rendering succeeds even when brand/ is absent — no hard failure."""
        src = self.tmp / "doc.md"
        src.write_text("# Hello\n")
        bin_dir = self.tmp / "bin"
        bin_dir.mkdir()
        self._make_fake_pandoc(bin_dir)
        result = run(
            ["bash", str(ORCH), str(src)],
            cwd=self.tmp,
            env={**os.environ, "PATH": f"{bin_dir}:{os.environ.get('PATH', '')}"},
        )
        self.assertEqual(result.returncode, 0)
        self.assertTrue((self.tmp / "doc.docx").exists())


class TestGenerateTemplatesPartialDeps(unittest.TestCase):
    """Covers exit-code behavior of generate-templates.py when only some deps
    are missing. Regression test for #275: the script must exit 0 when at
    least one format succeeded, and only exit 1 when ALL formats failed.
    """

    def setUp(self) -> None:
        # Require at least docx and openpyxl so DOCX + XLSX can succeed.
        for dep in ("docx", "openpyxl"):
            try:
                __import__(dep)
            except ImportError:
                self.skipTest(f"'{dep}' not installed; skipping partial-dep test")
        self.tmp = Path(tempfile.mkdtemp(prefix="bsg-gt-"))
        brand = self.tmp / "brand"
        brand.mkdir()
        tokens = {"name": "Test", "colors": {"primary": "#1a1a2e"}, "fonts": {"primary": "Helvetica Neue"}}
        (brand / "tokens.json").write_text(json.dumps(tokens))
        # Helper wrapper that blocks pptx import then runs generate-templates.py
        self.wrapper = self.tmp / "_block_pptx.py"
        self.wrapper.write_text(
            "import builtins, runpy, sys\n"
            "_real = builtins.__import__\n"
            "def _fake(name, *args, **kwargs):\n"
            "    if name == 'pptx' or name.startswith('pptx.'): raise ImportError(f'No module named {name!r}')\n"
            "    return _real(name, *args, **kwargs)\n"
            "builtins.__import__ = _fake\n"
            "runpy.run_path(sys.argv[1], run_name='__main__')\n"
        )

    def tearDown(self) -> None:
        shutil.rmtree(self.tmp, ignore_errors=True)

    def _run_with_pptx_blocked(self) -> subprocess.CompletedProcess[str]:
        return run(
            ["python3", str(self.wrapper), str(GEN_TEMPLATES)],
            cwd=self.tmp,
            check=False,
        )

    def test_exits_zero_when_only_pptx_dep_missing(self) -> None:
        """#275: partial success (docx+xlsx ok, pptx skipped) must exit 0."""
        result = self._run_with_pptx_blocked()
        self.assertEqual(
            result.returncode, 0,
            f"Expected exit 0 when only pptx dep missing, got {result.returncode}.\n"
            f"stderr={result.stderr}",
        )

    def test_summary_line_printed_on_partial_success(self) -> None:
        """#275: a one-line summary must be printed showing per-format status."""
        result = self._run_with_pptx_blocked()
        combined = result.stdout + result.stderr
        self.assertIn("templates:", combined.lower())
        self.assertIn("pptx", combined.lower())

    def test_exits_nonzero_when_all_formats_fail(self) -> None:
        """If every format fails (bad tokens), exit 1 is still correct."""
        # Corrupt tokens so all generators raise.
        (self.tmp / "brand" / "tokens.json").write_text("{}")
        result = run(
            ["python3", str(GEN_TEMPLATES)],
            cwd=self.tmp,
            check=False,
        )
        self.assertNotEqual(result.returncode, 0)


class TestInitBrandSmoke(unittest.TestCase):
    """End-to-end init: scan + generate templates. Skipped when deps missing."""

    REQUIRED_DEPS = ("docx", "pptx", "openpyxl")

    def setUp(self) -> None:
        for dep in self.REQUIRED_DEPS:
            try:
                __import__(dep)
            except ImportError:
                self.skipTest(
                    f"'{dep}' not installed; run scripts/install-local.sh to enable init tests"
                )
        if not shutil.which("pandoc"):
            self.skipTest("pandoc not installed; run scripts/install-local.sh")
        self.tmp = Path(tempfile.mkdtemp(prefix="bsg-init-"))
        (self.tmp / "README.md").write_text("# Init Test Repo\n\nDoes things.\n")
        css = self.tmp / "styles.css"
        css.write_text(":root { --primary-color: #3a86ff; --font-family-primary: 'Roboto'; }\n")

    def tearDown(self) -> None:
        if hasattr(self, "tmp"):
            shutil.rmtree(self.tmp, ignore_errors=True)

    def test_init_creates_tokens_json(self) -> None:
        run(["bash", str(INIT_BRAND)], cwd=self.tmp)
        tokens_path = self.tmp / "brand" / "tokens.json"
        self.assertTrue(tokens_path.exists())
        tokens = json.loads(tokens_path.read_text())
        self.assertEqual(tokens["name"], "Init Test Repo")
        self.assertEqual(tokens["colors"]["primary"], "#3a86ff")
        self.assertIn("Roboto", tokens["fonts"]["primary"])

    def test_init_creates_all_three_templates(self) -> None:
        run(["bash", str(INIT_BRAND)], cwd=self.tmp)
        templates = self.tmp / "brand" / "templates"
        for fname in ("reference.docx", "template.pptx", "template.xlsx"):
            path = templates / fname
            self.assertTrue(path.exists(), f"Missing {fname}")
            self.assertTrue(path.stat().st_size > 0)
            self.assertTrue(path.read_bytes().startswith(ZIP_MAGIC), f"{fname} is not a ZIP")

    def test_init_is_idempotent_without_force(self) -> None:
        run(["bash", str(INIT_BRAND)], cwd=self.tmp)
        # Second run without --force: should exit 0, not overwrite.
        result = run(["bash", str(INIT_BRAND)], cwd=self.tmp)
        self.assertIn("already exist", result.stdout)

    def test_init_force_overwrites(self) -> None:
        run(["bash", str(INIT_BRAND)], cwd=self.tmp)
        # Modify tokens, re-run with --force — templates should be regenerated.
        tokens_path = self.tmp / "brand" / "tokens.json"
        tokens = json.loads(tokens_path.read_text())
        tokens["colors"]["primary"] = "#ff0000"
        tokens_path.write_text(json.dumps(tokens))
        run(["bash", str(INIT_BRAND), "--force"], cwd=self.tmp)
        # If it succeeded, templates exist and are fresh.
        self.assertTrue((self.tmp / "brand" / "templates" / "reference.docx").exists())

    def test_tokens_only_skips_template_generation(self) -> None:
        run(["bash", str(INIT_BRAND), "--tokens-only"], cwd=self.tmp)
        self.assertTrue((self.tmp / "brand" / "tokens.json").exists())
        # tokens-only should not create Office templates, but brand-audit.md is OK
        for fname in ("reference.docx", "template.pptx", "template.xlsx"):
            self.assertFalse(
                (self.tmp / "brand" / "templates" / fname).exists(),
                f"{fname} should not exist with --tokens-only",
            )

    def test_dry_run_writes_tokens_and_audit_only(self) -> None:
        run(["bash", str(INIT_BRAND), "--dry-run"], cwd=self.tmp)
        self.assertTrue((self.tmp / "brand" / "tokens.json").exists())
        self.assertTrue((self.tmp / "brand" / "templates" / "brand-audit.md").exists())
        for fname in ("reference.docx", "template.pptx", "template.xlsx"):
            self.assertFalse(
                (self.tmp / "brand" / "templates" / fname).exists(),
                f"{fname} should not exist with --dry-run",
            )

    def test_dry_run_skips_idempotency_guard(self) -> None:
        """--dry-run should work even when templates already exist."""
        run(["bash", str(INIT_BRAND)], cwd=self.tmp)
        # Without --force, a second full run would exit early.
        # But --dry-run should still scan and produce the audit.
        result = run(["bash", str(INIT_BRAND), "--dry-run"], cwd=self.tmp)
        self.assertIn("Dry run", result.stdout)
        self.assertTrue((self.tmp / "brand" / "templates" / "brand-audit.md").exists())

    def test_init_creates_brand_audit_md(self) -> None:
        run(["bash", str(INIT_BRAND)], cwd=self.tmp)
        audit = self.tmp / "brand" / "templates" / "brand-audit.md"
        self.assertTrue(audit.exists())
        content = audit.read_text()
        self.assertIn("Brand Audit", content)
        self.assertIn("Init Test Repo", content)


class TestRenderDocxSmoke(unittest.TestCase):
    """End-to-end DOCX render via the real pandoc.

    Skipped when pandoc isn't installed — devs run install-local.sh to
    opt in. In CI, pandoc is a prerequisite of this skill's test job.
    """

    def setUp(self) -> None:
        if not shutil.which("pandoc"):
            self.skipTest("pandoc not installed; run scripts/install-local.sh")
        self.tmp = Path(tempfile.mkdtemp(prefix="bsg-smoke-"))
        self.src = self.tmp / "report.md"
        self.src.write_text("# Title\n\nParagraph with **bold** text.\n")

    def tearDown(self) -> None:
        if hasattr(self, "tmp"):
            shutil.rmtree(self.tmp, ignore_errors=True)

    def test_renders_unbranded_docx_without_template(self) -> None:
        """No brand/templates/ — render-docx.sh still succeeds silently.

        The 'unbranded' onboarding message is now the orchestrator's job
        (md-to-office.sh). render-docx.sh is a pure renderer: it just
        omits --reference-doc when no template is given, with no noise.
        """
        out = self.tmp / "report.docx"
        result = run(
            ["bash", str(RENDER_DOCX), str(self.src), str(out)],
            cwd=self.tmp,
        )
        self.assertTrue(out.exists())
        self.assertTrue(out.stat().st_size > 0)
        self.assertTrue(out.read_bytes().startswith(ZIP_MAGIC))
        self.assertEqual(result.returncode, 0)

    def test_renders_end_to_end_via_orchestrator(self) -> None:
        out = self.tmp / "report.docx"
        run(["bash", str(ORCH), str(self.src)], cwd=self.tmp)
        self.assertTrue(out.exists())
        self.assertTrue(out.read_bytes().startswith(ZIP_MAGIC))


class TestMarkdownParser(unittest.TestCase):
    """Covers the markdown-to-slides parser in render-pptx.py. Pure Python, no deps."""

    def setUp(self) -> None:
        import importlib.util
        spec = importlib.util.spec_from_file_location("render_pptx", str(RENDER_PPTX_PY))
        mod = importlib.util.module_from_spec(spec)
        sys.modules["render_pptx"] = mod
        spec.loader.exec_module(mod)
        self.mod = mod

    def test_h1_creates_title_slide(self) -> None:
        slides = self.mod.parse_markdown("# Hello World\n")
        self.assertEqual(len(slides), 1)
        self.assertEqual(slides[0].layout, "title")
        self.assertEqual(slides[0].title, "Hello World")

    def test_h2_creates_content_slide(self) -> None:
        slides = self.mod.parse_markdown("## Section One\n\nBody text.\n")
        self.assertEqual(len(slides), 1)
        self.assertEqual(slides[0].layout, "content")
        self.assertEqual(slides[0].title, "Section One")

    def test_multiple_h2_create_multiple_slides(self) -> None:
        md = "## Slide 1\n\nText.\n\n## Slide 2\n\nMore text.\n"
        slides = self.mod.parse_markdown(md)
        self.assertEqual(len(slides), 2)

    def test_horizontal_rule_creates_slide_break(self) -> None:
        md = "## Slide A\n\nContent.\n\n---\n\n## Slide B\n\nContent.\n"
        slides = self.mod.parse_markdown(md)
        self.assertEqual(len(slides), 2)

    def test_bullets_parsed(self) -> None:
        md = "## List Slide\n\n- Item one\n- Item two\n"
        slides = self.mod.parse_markdown(md)
        self.assertEqual(len(slides), 1)
        bullets_block = [b for b in slides[0].blocks if b[0] == "bullets"]
        self.assertEqual(len(bullets_block), 1)
        self.assertEqual(len(bullets_block[0][1]), 2)

    def test_nested_bullets_have_levels(self) -> None:
        md = "## Nested\n\n- Top\n  - Nested\n"
        slides = self.mod.parse_markdown(md)
        bullets = [b for b in slides[0].blocks if b[0] == "bullets"][0][1]
        self.assertEqual(bullets[0][0], 0)  # level 0
        self.assertEqual(bullets[1][0], 1)  # level 1

    def test_table_parsed(self) -> None:
        md = "## Data\n\n| A | B |\n|---|---|\n| 1 | 2 |\n"
        slides = self.mod.parse_markdown(md)
        table_block = [b for b in slides[0].blocks if b[0] == "table"]
        self.assertEqual(len(table_block), 1)
        self.assertEqual(table_block[0][1]["headers"], ["A", "B"])
        self.assertEqual(len(table_block[0][1]["rows"]), 1)

    def test_image_parsed(self) -> None:
        md = "## Pic\n\n![logo](img/logo.png)\n"
        slides = self.mod.parse_markdown(md)
        img_block = [b for b in slides[0].blocks if b[0] == "image"]
        self.assertEqual(len(img_block), 1)
        self.assertEqual(img_block[0][1]["path"], "img/logo.png")

    def test_code_block_parsed(self) -> None:
        md = "## Code\n\n```python\nprint('hi')\n```\n"
        slides = self.mod.parse_markdown(md)
        code_block = [b for b in slides[0].blocks if b[0] == "code"]
        self.assertEqual(len(code_block), 1)
        self.assertIn("print", code_block[0][1])

    def test_frontmatter_stripped(self) -> None:
        md = "---\ngamma:\n  format: presentation\n---\n# Title\n"
        body = self.mod.strip_frontmatter(md)
        self.assertNotIn("gamma", body)
        self.assertIn("# Title", body)

    def test_h3_stays_in_current_slide(self) -> None:
        md = "## Main\n\n### Sub\n\nText.\n"
        slides = self.mod.parse_markdown(md)
        self.assertEqual(len(slides), 1)
        headings = [b for b in slides[0].blocks if b[0] == "heading"]
        self.assertEqual(len(headings), 1)
        self.assertEqual(headings[0][1], "Sub")

    def test_mixed_content_slide(self) -> None:
        md = (
            "# Title\n\nSubtitle.\n\n## Content\n\n"
            "- Bullet\n\n| A |\n|---|\n| 1 |\n\nParagraph.\n"
        )
        slides = self.mod.parse_markdown(md)
        self.assertEqual(len(slides), 2)
        self.assertEqual(slides[0].layout, "title")
        self.assertEqual(slides[1].layout, "content")
        block_types = [b[0] for b in slides[1].blocks]
        self.assertIn("bullets", block_types)
        self.assertIn("table", block_types)
        self.assertIn("text", block_types)


class TestRenderPptxSmoke(unittest.TestCase):
    """End-to-end PPTX render. Skipped when python-pptx isn't installed."""

    def setUp(self) -> None:
        try:
            __import__("pptx")
        except ImportError:
            self.skipTest("python-pptx not installed; run scripts/install-local.sh")
        self.tmp = Path(tempfile.mkdtemp(prefix="bsg-pptx-"))
        self.src = self.tmp / "deck.md"
        self.src.write_text(
            "# Presentation Title\n\nSubtitle here.\n\n"
            "## Slide One\n\n- Bullet **one**\n- Bullet *two*\n\n"
            "## Slide Two\n\n| Col A | Col B |\n|-------|-------|\n| val 1 | val 2 |\n\n"
            "## Slide Three\n\n```python\nprint('hello')\n```\n"
        )

    def tearDown(self) -> None:
        if hasattr(self, "tmp"):
            shutil.rmtree(self.tmp, ignore_errors=True)

    def test_renders_unbranded_pptx(self) -> None:
        out = self.tmp / "deck.pptx"
        result = run(
            ["bash", str(RENDER_PPTX), str(self.src), str(out)],
            cwd=self.tmp,
        )
        self.assertTrue(out.exists())
        self.assertTrue(out.stat().st_size > 0)
        self.assertTrue(out.read_bytes().startswith(ZIP_MAGIC))
        self.assertEqual(result.returncode, 0)

    def test_pptx_has_correct_slide_count(self) -> None:
        out = self.tmp / "deck.pptx"
        run(["bash", str(RENDER_PPTX), str(self.src), str(out)], cwd=self.tmp)
        from pptx import Presentation
        prs = Presentation(str(out))
        self.assertEqual(len(prs.slides), 4)

    def test_renders_end_to_end_via_orchestrator(self) -> None:
        out = self.tmp / "deck.pptx"
        run(
            ["bash", str(ORCH), str(self.src), "--target", "pptx"],
            cwd=self.tmp,
        )
        self.assertTrue(out.exists())
        self.assertTrue(out.read_bytes().startswith(ZIP_MAGIC))

    def test_autodetect_via_orchestrator(self) -> None:
        src = self.tmp / "auto.md"
        src.write_text("---\ngamma:\n  format: presentation\n---\n# Auto\n\nBody.\n")
        run(["bash", str(ORCH), str(src)], cwd=self.tmp)
        self.assertTrue((self.tmp / "auto.pptx").exists())

    def test_renders_with_template(self) -> None:
        from pptx import Presentation
        tmpl_dir = self.tmp / "brand" / "templates"
        tmpl_dir.mkdir(parents=True)
        prs = Presentation()
        tmpl_path = tmpl_dir / "template.pptx"
        prs.save(str(tmpl_path))
        out = self.tmp / "deck.pptx"
        run(
            ["bash", str(RENDER_PPTX), str(self.src), str(out),
             "--template", str(tmpl_path)],
            cwd=self.tmp,
        )
        self.assertTrue(out.exists())
        self.assertTrue(out.stat().st_size > 0)


if __name__ == "__main__":
    unittest.main(verbosity=2)
