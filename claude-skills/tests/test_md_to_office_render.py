#!/usr/bin/env python3
"""Smoke tests for render-docx.sh, render-pptx.sh, and the markdown parser.

Split out of test_md_to_office.py for #330 (single file >500 LOC).
Run individually:

    python3 claude-skills/tests/test_md_to_office_render.py

Stdlib only; auto-skips when external deps are missing.
"""

from __future__ import annotations

import shutil
import sys
import tempfile
import unittest
from pathlib import Path

from _md_to_office_helpers import (
    ORCH,
    RENDER_DOCX,
    RENDER_PPTX,
    RENDER_PPTX_PY,
    ZIP_MAGIC,
    run,
)


class TestRenderDocxSmoke(unittest.TestCase):
    """End-to-end DOCX render via the real pandoc.

    Skipped when pandoc isn't installed — devs run install-local.sh to
    opt in. In CI, pandoc is a prerequisite of this skill's test job.
    """

    def setUp(self) -> None:
        if not shutil.which("pandoc"):
            self.skipTest("pandoc not installed; run scripts/install-local.sh")
        self.tmp = Path(tempfile.mkdtemp(prefix="bsg-smoke-"))
        self.src = self.tmp / "report.md"
        self.src.write_text("# Title\n\nParagraph with **bold** text.\n")

    def tearDown(self) -> None:
        if hasattr(self, "tmp"):
            shutil.rmtree(self.tmp, ignore_errors=True)

    def test_renders_unbranded_docx_without_template(self) -> None:
        """No brand/templates/ — render-docx.sh still succeeds silently.

        The 'unbranded' onboarding message is now the orchestrator's job
        (md-to-office.sh). render-docx.sh is a pure renderer: it just
        omits --reference-doc when no template is given, with no noise.
        """
        out = self.tmp / "report.docx"
        result = run(
            ["bash", str(RENDER_DOCX), str(self.src), str(out)],
            cwd=self.tmp,
        )
        self.assertTrue(out.exists())
        self.assertTrue(out.stat().st_size > 0)
        self.assertTrue(out.read_bytes().startswith(ZIP_MAGIC))
        self.assertEqual(result.returncode, 0)

    def test_renders_end_to_end_via_orchestrator(self) -> None:
        out = self.tmp / "report.docx"
        run(["bash", str(ORCH), str(self.src)], cwd=self.tmp)
        self.assertTrue(out.exists())
        self.assertTrue(out.read_bytes().startswith(ZIP_MAGIC))

class TestMarkdownParser(unittest.TestCase):
    """Covers the markdown-to-slides parser in render-pptx.py. Pure Python, no deps."""

    def setUp(self) -> None:
        import importlib.util
        spec = importlib.util.spec_from_file_location("render_pptx", str(RENDER_PPTX_PY))
        mod = importlib.util.module_from_spec(spec)
        sys.modules["render_pptx"] = mod
        spec.loader.exec_module(mod)
        self.mod = mod

    def test_h1_creates_title_slide(self) -> None:
        slides = self.mod.parse_markdown("# Hello World\n")
        self.assertEqual(len(slides), 1)
        self.assertEqual(slides[0].layout, "title")
        self.assertEqual(slides[0].title, "Hello World")

    def test_h2_creates_content_slide(self) -> None:
        slides = self.mod.parse_markdown("## Section One\n\nBody text.\n")
        self.assertEqual(len(slides), 1)
        self.assertEqual(slides[0].layout, "content")
        self.assertEqual(slides[0].title, "Section One")

    def test_multiple_h2_create_multiple_slides(self) -> None:
        md = "## Slide 1\n\nText.\n\n## Slide 2\n\nMore text.\n"
        slides = self.mod.parse_markdown(md)
        self.assertEqual(len(slides), 2)

    def test_horizontal_rule_creates_slide_break(self) -> None:
        md = "## Slide A\n\nContent.\n\n---\n\n## Slide B\n\nContent.\n"
        slides = self.mod.parse_markdown(md)
        self.assertEqual(len(slides), 2)

    def test_bullets_parsed(self) -> None:
        md = "## List Slide\n\n- Item one\n- Item two\n"
        slides = self.mod.parse_markdown(md)
        self.assertEqual(len(slides), 1)
        bullets_block = [b for b in slides[0].blocks if b[0] == "bullets"]
        self.assertEqual(len(bullets_block), 1)
        self.assertEqual(len(bullets_block[0][1]), 2)

    def test_nested_bullets_have_levels(self) -> None:
        md = "## Nested\n\n- Top\n  - Nested\n"
        slides = self.mod.parse_markdown(md)
        bullets = [b for b in slides[0].blocks if b[0] == "bullets"][0][1]
        self.assertEqual(bullets[0][0], 0)  # level 0
        self.assertEqual(bullets[1][0], 1)  # level 1

    def test_table_parsed(self) -> None:
        md = "## Data\n\n| A | B |\n|---|---|\n| 1 | 2 |\n"
        slides = self.mod.parse_markdown(md)
        table_block = [b for b in slides[0].blocks if b[0] == "table"]
        self.assertEqual(len(table_block), 1)
        self.assertEqual(table_block[0][1]["headers"], ["A", "B"])
        self.assertEqual(len(table_block[0][1]["rows"]), 1)

    def test_image_parsed(self) -> None:
        md = "## Pic\n\n![logo](img/logo.png)\n"
        slides = self.mod.parse_markdown(md)
        img_block = [b for b in slides[0].blocks if b[0] == "image"]
        self.assertEqual(len(img_block), 1)
        self.assertEqual(img_block[0][1]["path"], "img/logo.png")

    def test_code_block_parsed(self) -> None:
        md = "## Code\n\n```python\nprint('hi')\n```\n"
        slides = self.mod.parse_markdown(md)
        code_block = [b for b in slides[0].blocks if b[0] == "code"]
        self.assertEqual(len(code_block), 1)
        self.assertIn("print", code_block[0][1])

    def test_frontmatter_stripped(self) -> None:
        md = "---\ngamma:\n  format: presentation\n---\n# Title\n"
        body = self.mod.strip_frontmatter(md)
        self.assertNotIn("gamma", body)
        self.assertIn("# Title", body)

    def test_h3_stays_in_current_slide(self) -> None:
        md = "## Main\n\n### Sub\n\nText.\n"
        slides = self.mod.parse_markdown(md)
        self.assertEqual(len(slides), 1)
        headings = [b for b in slides[0].blocks if b[0] == "heading"]
        self.assertEqual(len(headings), 1)
        self.assertEqual(headings[0][1], "Sub")

    def test_mixed_content_slide(self) -> None:
        md = (
            "# Title\n\nSubtitle.\n\n## Content\n\n"
            "- Bullet\n\n| A |\n|---|\n| 1 |\n\nParagraph.\n"
        )
        slides = self.mod.parse_markdown(md)
        self.assertEqual(len(slides), 2)
        self.assertEqual(slides[0].layout, "title")
        self.assertEqual(slides[1].layout, "content")
        block_types = [b[0] for b in slides[1].blocks]
        self.assertIn("bullets", block_types)
        self.assertIn("table", block_types)
        self.assertIn("text", block_types)

class TestRenderPptxSmoke(unittest.TestCase):
    """End-to-end PPTX render. Skipped when python-pptx isn't installed."""

    def setUp(self) -> None:
        try:
            __import__("pptx")
        except ImportError:
            self.skipTest("python-pptx not installed; run scripts/install-local.sh")
        self.tmp = Path(tempfile.mkdtemp(prefix="bsg-pptx-"))
        self.src = self.tmp / "deck.md"
        self.src.write_text(
            "# Presentation Title\n\nSubtitle here.\n\n"
            "## Slide One\n\n- Bullet **one**\n- Bullet *two*\n\n"
            "## Slide Two\n\n| Col A | Col B |\n|-------|-------|\n| val 1 | val 2 |\n\n"
            "## Slide Three\n\n```python\nprint('hello')\n```\n"
        )

    def tearDown(self) -> None:
        if hasattr(self, "tmp"):
            shutil.rmtree(self.tmp, ignore_errors=True)

    def test_renders_unbranded_pptx(self) -> None:
        out = self.tmp / "deck.pptx"
        result = run(
            ["bash", str(RENDER_PPTX), str(self.src), str(out)],
            cwd=self.tmp,
        )
        self.assertTrue(out.exists())
        self.assertTrue(out.stat().st_size > 0)
        self.assertTrue(out.read_bytes().startswith(ZIP_MAGIC))
        self.assertEqual(result.returncode, 0)

    def test_pptx_has_correct_slide_count(self) -> None:
        out = self.tmp / "deck.pptx"
        run(["bash", str(RENDER_PPTX), str(self.src), str(out)], cwd=self.tmp)
        from pptx import Presentation
        prs = Presentation(str(out))
        self.assertEqual(len(prs.slides), 4)

    def test_renders_end_to_end_via_orchestrator(self) -> None:
        out = self.tmp / "deck.pptx"
        run(
            ["bash", str(ORCH), str(self.src), "--target", "pptx"],
            cwd=self.tmp,
        )
        self.assertTrue(out.exists())
        self.assertTrue(out.read_bytes().startswith(ZIP_MAGIC))

    def test_autodetect_via_orchestrator(self) -> None:
        src = self.tmp / "auto.md"
        src.write_text("---\ngamma:\n  format: presentation\n---\n# Auto\n\nBody.\n")
        run(["bash", str(ORCH), str(src)], cwd=self.tmp)
        self.assertTrue((self.tmp / "auto.pptx").exists())

    def test_renders_with_template(self) -> None:
        from pptx import Presentation
        tmpl_dir = self.tmp / "brand" / "templates"
        tmpl_dir.mkdir(parents=True)
        prs = Presentation()
        tmpl_path = tmpl_dir / "template.pptx"
        prs.save(str(tmpl_path))
        out = self.tmp / "deck.pptx"
        run(
            ["bash", str(RENDER_PPTX), str(self.src), str(out),
             "--template", str(tmpl_path)],
            cwd=self.tmp,
        )
        self.assertTrue(out.exists())
        self.assertTrue(out.stat().st_size > 0)


if __name__ == "__main__":
    unittest.main(verbosity=2)
